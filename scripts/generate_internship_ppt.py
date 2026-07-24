#!/usr/bin/env python3
"""Generate AegisSOC internship presentation PPTX for CE446.

Mirrors the structure/style of the DEPSTAR internship demo PPT while using
content and screenshots from the final report (23DCE081_report.docx).
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# --- Paths ---
REPORT = Path("/Users/harshil/Downloads/23DCE081_report.docx")
OUT = Path("/Users/harshil/Downloads/23DCE081_Internship_PPT.pptx")
MEDIA = Path("/tmp/aegis_report_media")

# Widescreen 16:9 (same as demo)
SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)

BLUE = RGBColor(0x44, 0x72, 0xC4)
DARK = RGBColor(0x1F, 0x1F, 0x1F)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xD6, 0xE3, 0xF8)
ACCENT_BAR = RGBColor(0x2F, 0x54, 0x96)

# Report media mapping (document.xml.rels)
FIG = {
    "5.1": MEDIA / "image4.png",  # End-to-end data flow
    "5.2": MEDIA / "image5.png",  # Microservices architecture
    "5.3": MEDIA / "image6.png",  # Alert-to-response sequence
    "6.1": MEDIA / "image7.png",  # Alert Queue
    "6.2": MEDIA / "image8.png",  # Cases
    "6.3": MEDIA / "image9.png",  # Case Detail
    "6.4": MEDIA / "image10.png",  # Investigation
    "6.5": MEDIA / "image11.png",  # Response
    "6.6": MEDIA / "image12.png",  # Approval modal
    "6.7": MEDIA / "image13.png",  # Audit
    "6.8": MEDIA / "image14.png",  # Metrics
    "6.9": MEDIA / "image15.png",  # Replay (legacy)
}


def ensure_media() -> None:
    MEDIA.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(REPORT) as z:
        for name in z.namelist():
            if name.startswith("word/media/"):
                dest = MEDIA / Path(name).name
                if not dest.exists():
                    dest.write_bytes(z.read(name))


def set_run(run, *, size=18, bold=False, color=DARK, name="Times New Roman"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("a:latin"))
    # force East Asian / latin
    from pptx.oxml import parse_xml

    # python-pptx sets latin via font.name; also set for compatibility
    run.font.name = name


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Times New Roman"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=color, name=font)
    return box


def add_bullets(slide, left, top, width, height, items, *, size=16, color=DARK, bold_first=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = f"•  {item}"
        set_run(run, size=size, bold=(bold_first and i == 0), color=color)
    return box


def add_bar(slide, top=True):
    from pptx.enum.shapes import MSO_SHAPE

    if top:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Emu(278813))
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Emu(278813), SLIDE_W, Emu(278813))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BAR
    shape.line.fill.background()
    return shape


def add_title_banner(slide, title: str):
    from pptx.enum.shapes import MSO_SHAPE

    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.55))
    banner.fill.solid()
    banner.fill.fore_color.rgb = BLUE
    banner.line.fill.background()
    tf = banner.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    set_run(run, size=22, bold=True, color=WHITE)
    return banner


def fit_picture(slide, path: Path, left, top, max_w, max_h):
    from PIL import Image

    im = Image.open(path)
    iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    # center horizontally in the max box
    x = int(left + (max_w - w) / 2)
    y = int(top + (max_h - h) / 2)
    return slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_caption_chip(slide, text: str):
    from pptx.enum.shapes import MSO_SHAPE

    chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.15), Inches(3.0), Inches(0.32))
    chip.fill.solid()
    chip.fill.fore_color.rgb = BLUE
    chip.line.fill.background()
    tf = chip.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size=12, bold=True, color=WHITE)
    return chip


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# --------------------------------------------------------------------------- slides


def slide_title(prs: Presentation):
    s = blank_slide(prs)
    add_bar(s, top=True)
    add_bar(s, top=False)

    add_textbox(
        s,
        Inches(0.6),
        Inches(0.45),
        Inches(8.8),
        Inches(1.1),
        "CHAROTAR UNIVERSITY OF SCIENCE & TECHNOLOGY (CHARUSAT)\n"
        "FACULTY OF TECHNOLOGY & ENGINEERING\n"
        "DEVANG PATEL INSTITUTE OF ADVANCE TECHNOLOGY & RESEARCH (DEPSTAR)",
        size=14,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )

    # divider
    from pptx.enum.shapes import MSO_SHAPE

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.55), Inches(8.4), Emu(25000))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    add_textbox(
        s,
        Inches(1.5),
        Inches(1.7),
        Inches(7.0),
        Inches(1.1),
        "Department of Computer Engineering\n"
        "SUMMER INTERNSHIP-II [CE446]\n"
        "Summer Internship Project",
        size=16,
        bold=False,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        s,
        Inches(1.5),
        Inches(2.85),
        Inches(7.0),
        Inches(0.7),
        "AegisSOC: AI-Assisted SOC Triage Platform",
        size=22,
        bold=True,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        s,
        Inches(0.5),
        Inches(3.9),
        Inches(4.5),
        Inches(1.2),
        "Name of Student & Id:\nHarshil Patel — 23DCE081",
        size=14,
        bold=False,
        color=DARK,
        align=PP_ALIGN.LEFT,
    )
    add_textbox(
        s,
        Inches(5.2),
        Inches(3.9),
        Inches(4.3),
        Inches(1.3),
        "Internal Guide: Dr. Parth Goel\n"
        "External Guide: Nikunj Bhavsar\n"
        "Human Resource Executive, Digiesale",
        size=14,
        bold=False,
        color=DARK,
        align=PP_ALIGN.LEFT,
    )


def slide_problem(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "AegisSOC — Problem Statement")
    add_bullets(
        s,
        Inches(0.7),
        Inches(1.2),
        Inches(8.6),
        Inches(4.0),
        [
            "SOC analysts face alert fatigue from high false-positive rates across fragmented tools.",
            "Events stay siloed across endpoint, network, identity, cloud, and email sources — context is lost.",
            "Manual correlation slows investigation; attack narratives are hard to reconstruct quickly.",
            "Disruptive response actions risk business impact without clear human approval and audit trails.",
            "AegisSOC addresses this with graph-aware correlation, explainable detection, and LLM-assisted (not LLM-driven) triage.",
        ],
        size=17,
    )


def slide_scope(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "Project Scope and Features")

    # three columns
    headers = ["Core Features", "Out of Scope / Limits", "Future Work"]
    cols = [
        [
            "JWT analyst console (alerts, cases, investigation).",
            "Ingest → normalize → enrich → graph → detect → case → triage → response pipeline.",
            "Temporal entity graph + attack-path visualization.",
            "Sigma-like rules, correlation, graph risk scoring.",
            "Human-in-the-loop approvals with audit logging.",
            "Docker Compose + sync-mode demo scenarios.",
        ],
        [
            "Not a full commercial SIEM replacement.",
            "Live VirusTotal/OpenCTI are optional stubs by default.",
            "Graph ML is lightweight (numpy), not GPU GNN.",
            "Tenant isolation is logical (tenant_id), not physical.",
        ],
        [
            "Multi-broker Redpanda + autoscaling.",
            "Neo4j cluster / AuraDB replicas.",
            "Live threat-intel feeds with caching.",
            "Heterogeneous GNN risk scoring.",
            "mTLS + per-source API keys at ingestion.",
        ],
    ]
    lefts = [0.4, 3.55, 6.7]
    from pptx.enum.shapes import MSO_SHAPE

    for i, (h, items) in enumerate(zip(headers, cols)):
        x = Inches(lefts[i])
        head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.05), Inches(2.9), Inches(0.4))
        head.fill.solid()
        head.fill.fore_color.rgb = BLUE
        head.line.fill.background()
        tf = head.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        set_run(run, size=13, bold=True, color=WHITE)
        add_bullets(s, x, Inches(1.55), Inches(2.9), Inches(3.5), items, size=12)


def slide_tech(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "Technology Stack and Tools")

    items = [
        ("1", "Backend", "Python 3.11, FastAPI (11 microservices)"),
        ("2", "Frontend", "React, TypeScript, Vite, Cytoscape.js"),
        ("3", "Data stores", "PostgreSQL, Neo4j, Redis, OpenSearch"),
        ("4", "Streaming", "Redpanda (Kafka-compatible)"),
        ("5", "Auth & Ops", "JWT RBAC, Docker Compose, Prometheus"),
        ("6", "Other tools", "Git/GitHub, VS Code, evaluation scripts"),
    ]
    from pptx.enum.shapes import MSO_SHAPE

    positions = [
        (0.6, 1.2),
        (5.1, 1.2),
        (0.6, 2.35),
        (5.1, 2.35),
        (0.6, 3.5),
        (5.1, 3.5),
    ]
    for (num, title, desc), (x, y) in zip(items, positions):
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.45), Inches(0.45))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        tf = circ.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_run(run, size=14, bold=True, color=WHITE)
        add_textbox(s, Inches(x + 0.55), Inches(y - 0.05), Inches(3.6), Inches(0.55), f"{title}\n{desc}", size=14, bold=False)


def slide_architecture(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "Implementation Architecture")
    add_bullets(
        s,
        Inches(0.5),
        Inches(1.0),
        Inches(4.2),
        Inches(3.8),
        [
            "Event-driven microservices over Kafka/Redpanda topics.",
            "CanonicalEvent schema unifies Sysmon, Zeek, EDR, AD, CloudTrail, K8s, email.",
            "Graph builder upserts entities/edges into Neo4j (or in-memory sync store).",
            "Detection → case clustering → LLM triage with citation grounding.",
            "Frontend gateway (BFF) exposes a unified JWT API to the React console.",
            "Disruptive actions require senior/analyst approval (human-in-the-loop).",
        ],
        size=14,
    )
    if FIG["5.2"].exists():
        fit_picture(s, FIG["5.2"], Inches(4.8), Inches(1.0), Inches(4.7), Inches(3.9))


def slide_flowchart(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "Flow Chart — End-to-End Pipeline")
    if FIG["5.1"].exists():
        fit_picture(s, FIG["5.1"], Inches(0.6), Inches(1.0), Inches(8.8), Inches(3.7))
    add_caption_chip(s, "Fig. 5.1  Data Flow")


def slide_sequence(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "Alert-to-Response Sequence")
    if FIG["5.3"].exists():
        fit_picture(s, FIG["5.3"], Inches(0.4), Inches(0.95), Inches(9.2), Inches(3.9))
    add_caption_chip(s, "Fig. 5.3  Sequence")


def slide_major(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "Major Functionalities, Limitations & Outcome")

    from pptx.enum.shapes import MSO_SHAPE

    blocks = [
        (
            "Major Functionalities",
            [
                "Multi-source ingest + scenario replay.",
                "Temporal graph & attack-path analysis.",
                "Explainable detection & case clustering.",
                "Grounded LLM triage + ATT&CK mapping.",
                "Approval workflow with full audit trail.",
            ],
        ),
        (
            "Limitations",
            [
                "Lightweight graph scorer (not GPU GNN).",
                "External intel connectors mostly stubbed.",
                "Heuristic prompt-injection filters.",
                "Logical multi-tenant scoping only.",
            ],
        ),
        (
            "Outcome",
            [
                "11 services + React console in Docker.",
                "Demo scenarios: ransomware, FP, graph memory.",
                "Detection stays deterministic; LLM assists.",
                "Human-approved response with evidence.",
            ],
        ),
    ]
    lefts = [0.35, 3.5, 6.65]
    for (title, items), x in zip(blocks, lefts):
        head = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.05), Inches(2.95), Inches(0.4))
        head.fill.solid()
        head.fill.fore_color.rgb = BLUE
        head.line.fill.background()
        tf = head.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        set_run(run, size=13, bold=True, color=WHITE)
        add_bullets(s, Inches(x), Inches(1.55), Inches(2.95), Inches(3.5), items, size=13)


def slide_screenshot(prs: Presentation, title: str, fig_key: str, caption: str):
    s = blank_slide(prs)
    add_title_banner(s, title)
    path = FIG[fig_key]
    if path.exists():
        fit_picture(s, path, Inches(0.5), Inches(0.95), Inches(9.0), Inches(4.0))
    add_caption_chip(s, caption)


def slide_references(prs: Presentation):
    s = blank_slide(prs)
    add_title_banner(s, "References and Resources")
    refs = [
        ("1", "MITRE ATT&CK Framework", "attack.mitre.org"),
        ("2", "Sigma Rule Format", "SigmaHQ / GitHub"),
        ("3", "FastAPI Documentation", "fastapi.tiangolo.com"),
        ("4", "Neo4j Graph Database", "neo4j.com/docs"),
        ("5", "Apache Kafka / Redpanda", "Streaming backbone"),
        ("6", "React + Docker Docs", "UI & deployment"),
    ]
    from pptx.enum.shapes import MSO_SHAPE

    positions = [(0.7, 1.2), (5.1, 1.2), (0.7, 2.35), (5.1, 2.35), (0.7, 3.5), (5.1, 3.5)]
    for (num, title, sub), (x, y) in zip(refs, positions):
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.45), Inches(0.45))
        circ.fill.solid()
        circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        tf = circ.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_run(run, size=14, bold=True, color=WHITE)
        add_textbox(s, Inches(x + 0.55), Inches(y), Inches(3.6), Inches(0.55), f"{title}\n{sub}", size=14)


def slide_thanks(prs: Presentation):
    s = blank_slide(prs)
    add_bar(s, top=True)
    add_bar(s, top=False)
    add_textbox(
        s,
        Inches(1.0),
        Inches(1.8),
        Inches(8.0),
        Inches(1.2),
        "Thank You !!!",
        size=40,
        bold=True,
        color=BLUE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        s,
        Inches(1.0),
        Inches(3.1),
        Inches(8.0),
        Inches(1.0),
        "Harshil Patel (23DCE081)\nAegisSOC — Digiesale | DEPSTAR, CHARUSAT",
        size=16,
        bold=False,
        color=DARK,
        align=PP_ALIGN.CENTER,
    )


def main() -> None:
    ensure_media()
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_scope(prs)
    slide_tech(prs)
    slide_architecture(prs)
    slide_flowchart(prs)
    slide_sequence(prs)
    slide_major(prs)

    # Screenshot slides (key UI from report Ch. 6)
    screenshots = [
        ("Screenshots — Alert Queue", "6.1", "Alert Queue"),
        ("Screenshots — Cases Management", "6.2", "Cases View"),
        ("Screenshots — Investigation Workspace", "6.4", "Attack-Path Graph"),
        ("Screenshots — Response & Approvals", "6.5", "Response Panel"),
        ("Screenshots — Approval Modal", "6.6", "Human Approval"),
        ("Screenshots — Metrics Dashboard", "6.8", "Ops Metrics"),
    ]
    for title, key, cap in screenshots:
        slide_screenshot(prs, title, key, cap)

    slide_references(prs)
    slide_thanks(prs)

    prs.save(OUT)
    # also copy into repo for backup
    repo_out = Path("/Users/harshil/Developer/AegisSOC/docs/23DCE081_Internship_PPT.pptx")
    repo_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(repo_out)
    print(f"Saved {OUT}")
    print(f"Saved {repo_out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
